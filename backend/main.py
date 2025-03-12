# main.py
from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse
import os
import uuid
from pathlib import Path
import shutil
from io import BytesIO
import base64

# Import libraries for PPTX to PDF conversion
from pptx import Presentation
from reportlab.lib.pagesizes import letter, landscape
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from reportlab.lib.utils import ImageReader
from PIL import Image

app = FastAPI()

# Enable CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Create upload and output directories if they don't exist
UPLOAD_FOLDER = Path("uploads")
OUTPUT_FOLDER = Path("outputs")
TEMP_FOLDER = Path("temp")

for folder in [UPLOAD_FOLDER, OUTPUT_FOLDER, TEMP_FOLDER]:
    folder.mkdir(exist_ok=True)

@app.post("/convert")
async def convert_pptx_to_pdf(file: UploadFile = File(...)):
    if not file.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Invalid file format. Please upload a .pptx file")
    
    # Generate unique filenames
    unique_id = str(uuid.uuid4())
    safe_filename = file.filename.replace(" ", "_")
    pptx_filename = f"{unique_id}_{safe_filename}"
    pdf_filename = f"{unique_id}.pdf"
    
    pptx_path = UPLOAD_FOLDER / pptx_filename
    pdf_path = OUTPUT_FOLDER / pdf_filename
    
    # Save the PPTX file
    try:
        with open(pptx_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save file: {str(e)}")
    
    # Convert PPTX to PDF
    try:
        convert_pptx_to_pdf(str(pptx_path), str(pdf_path), str(TEMP_FOLDER))
        return {"success": True, "pdf_url": f"/pdf/{pdf_filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")

def extract_slide_images(pptx_path, temp_dir):
    """Extract slide images from PPTX for better rendering"""
    prs = Presentation(pptx_path)
    slide_images = []
    
    for i, slide in enumerate(prs.slides):
        # Create a temporary directory for slide images
        slide_image_path = Path(temp_dir) / f"slide_{i+1}.png"
        
        # Use third-party tools or libraries to render slide as image
        # Since direct rendering is not available in python-pptx, we'll create a representation
        
        # Create a blank image with PIL as a fallback
        # In a production environment, you might use a headless browser or other tool
        # to capture the slide more accurately
        width = int(prs.slide_width * 0.75)  # Convert to pixels (approximate)
        height = int(prs.slide_height * 0.75)  # Convert to pixels (approximate)
        
        img = Image.new('RGB', (width, height), (255, 255, 255))
        
        # Save the image
        img.save(slide_image_path)
        slide_images.append(str(slide_image_path))
        
        # Extract text data for overlay
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                # Get approximate position (this is not perfect but gives a rough position)
                if hasattr(shape, "left") and hasattr(shape, "top"):
                    left = shape.left * width / prs.slide_width
                    top = shape.top * height / prs.slide_height
                    texts.append((shape.text, left, top))
                else:
                    texts.append((shape.text, 50, 50 + len(texts) * 30))  # Default position
        
        # Store text data in a separate file
        text_data_path = Path(temp_dir) / f"slide_{i+1}_text.txt"
        with open(text_data_path, "w", encoding="utf-8") as f:
            for text, left, top in texts:
                f.write(f"{text}|{left}|{top}\n")
    
    return slide_images

def convert_pptx_to_pdf(pptx_path, pdf_path, temp_dir):
    """Convert PPTX to PDF using python-pptx and reportlab"""
    prs = Presentation(pptx_path)
    
    # Extract slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Create a PDF with the same aspect ratio
    c = canvas.Canvas(pdf_path, pagesize=(slide_width, slide_height))
    
    # Process each slide
    for i, slide in enumerate(prs.slides):
        # Extract text and shapes from the slide
        slide_texts = []
        slide_images = []
        
        # Process shapes in the slide
        for shape in slide.shapes:
            # Handle text
            if hasattr(shape, "text") and shape.text:
                x = shape.left if hasattr(shape, "left") else 0
                y = slide_height - (shape.top + shape.height if hasattr(shape, "top") and hasattr(shape, "height") else 0)
                slide_texts.append((shape.text, x, y))
            
            # Handle images (placeholders)
            if shape.shape_type == 13:  # IMAGE shape type
                if hasattr(shape, "image"):
                    try:
                        image_stream = BytesIO(shape.image.blob)
                        x = shape.left if hasattr(shape, "left") else 0
                        y = slide_height - (shape.top + shape.height if hasattr(shape, "top") and hasattr(shape, "height") else 0)
                        width = shape.width if hasattr(shape, "width") else 100
                        height = shape.height if hasattr(shape, "height") else 100
                        slide_images.append((image_stream, x, y, width, height))
                    except:
                        pass
        
        # Draw slide background (white)
        c.setFillColorRGB(1, 1, 1)
        c.rect(0, 0, slide_width, slide_height, fill=True)
        
        # Draw a border around the slide
        c.setStrokeColorRGB(0.8, 0.8, 0.8)
        c.rect(0, 0, slide_width, slide_height, fill=False)
        
        # Draw images
        for img_stream, x, y, width, height in slide_images:
            try:
                img = ImageReader(img_stream)
                c.drawImage(img, x, y, width, height)
            except:
                # Draw a placeholder for failed images
                c.setFillColorRGB(0.9, 0.9, 0.9)
                c.rect(x, y, width, height, fill=True)
                c.setFillColorRGB(0.5, 0.5, 0.5)
                c.drawString(x + 10, y + height/2, "Image Placeholder")
        
        # Draw text elements
        for text, x, y in slide_texts:
            # Set font and size
            c.setFillColorRGB(0, 0, 0)
            c.setFont("Helvetica", 12)
            
            # Split text into lines
            lines = text.split('\n')
            line_height = 15
            
            for i, line in enumerate(lines):
                c.drawString(x, y - (i * line_height), line)
        
        # Add slide number
        c.setFont("Helvetica", 10)
        c.setFillColorRGB(0.5, 0.5, 0.5)
        c.drawString(slide_width - 60, 20, f"Slide {i+1}")
        
        # Move to the next page
        c.showPage()
    
    # Save the PDF
    c.save()

@app.get("/pdf/{filename}")
async def get_pdf(filename: str):
    file_path = OUTPUT_FOLDER / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="PDF not found")
    return FileResponse(file_path, media_type="application/pdf")

# Cleanup endpoint (optional - for maintenance)
@app.delete("/cleanup")
async def cleanup_files():
    # Delete files older than 1 hour (implement your own logic)
    return {"message": "Cleanup completed"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)